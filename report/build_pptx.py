#!/usr/bin/env python3
"""slides.html -> slides.pptx (full-bleed 2x screenshots + speaker notes).

A fallback for venues that require a slide upload. The HTML deck is the real
artefact — it is what gets presented; this exists so a conference upload form
or an unfamiliar presenting machine can never block the talk.

Each slide carries its section of SPEECH.md as PowerPoint speaker notes, so
presenter view shows the script beside the slide.
"""
import re, asyncio, io
from pathlib import Path

REPORT = Path("/home/chjin/lab_internship/appt_spgemm/report")
DECK   = REPORT / "slides.html"
SPEECH = REPORT / "SPEECH.md"
OUT    = REPORT / "slides.pptx"
SHOTS  = Path("/tmp/claude-1006/-home-chjin-lab-internship/"
              "d1470269-8c7d-458c-ae9d-9d67a44172e4/scratchpad/pptx_shots")

N_SLIDES = 8


def speaker_notes():
    """Per-slide script, markdown stripped, ready for presenter view."""
    raw = SPEECH.read_text(encoding="utf-8").split("## Handling Q&A")[0]
    out = {}
    for m in re.finditer(r"^## Slide (\d+) — (.+?)$(.*?)(?=^## |\Z)",
                         raw, re.M | re.S):
        n, title, body = int(m.group(1)), m.group(2), m.group(3)
        body = re.sub(r"^\*.*?\*$", "", body, flags=re.M)      # editorial notes
        body = re.sub(r"^---$", "", body, flags=re.M)
        body = body.replace("**", "").replace("`", "")
        body = re.sub(r"[ \t]+//", "   //", body)              # keep pause marks
        body = re.sub(r"\n{3,}", "\n\n", body).strip()
        out[n] = f"— {title} —\n\n{body}"
    return out


async def shoot():
    from playwright.async_api import async_playwright
    SHOTS.mkdir(parents=True, exist_ok=True)
    paths = []
    async with async_playwright() as p:
        b = await p.chromium.launch()
        pg = await b.new_page(viewport={"width": 1920, "height": 1080},
                              device_scale_factor=2)          # 3840x2160
        await pg.goto(DECK.as_uri())
        await pg.wait_for_timeout(700)
        for i in range(N_SLIDES):
            f = SHOTS / f"slide{i+1:02d}.png"
            await pg.screenshot(path=str(f))
            paths.append(f)
            if i < N_SLIDES - 1:
                await pg.keyboard.press("ArrowRight")
                await pg.wait_for_timeout(400)
        await b.close()
    return paths


def build(paths, notes):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)   # 16:9
    blank = prs.slide_layouts[6]
    for i, img in enumerate(paths, start=1):
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(img), 0, 0,
                             width=prs.slide_width, height=prs.slide_height)
        if i in notes:
            s.notes_slide.notes_text_frame.text = notes[i]
    prs.save(str(OUT))


notes = speaker_notes()
paths = asyncio.run(shoot())
build(paths, notes)
print(f"wrote {OUT}  ({len(paths)} slides, notes on {len(notes)})")
